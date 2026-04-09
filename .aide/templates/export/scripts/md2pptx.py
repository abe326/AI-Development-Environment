#!/usr/bin/env python3
"""
md2pptx.py — Marp Markdown → PowerPoint 変換スクリプト（v2）
aide-pm-slide で生成した Marp md から編集可能な PPTX を生成する。

【ナレッジベースで判明した既知問題への対応】
  - CJK フォント: font.name だけでは日本語に効かない。<a:ea>/<a:cs> をXML直接操作（Issue #768）
  - テキスト不可視: run.font.color.rgb を必ず明示。省略するとテーマ dk1 が自動適用される
  - 垂直中央: text_frame.vertical_anchor は信頼性低。bodyPr.set('anchor','ctr') で直接操作（Issue #191）
  - 先頭改行バグ: 先頭 \\n でフォントサイズが 18pt に固定される（Issue #740）→ lstrip で除去

使用方法:
  python3 md2pptx.py <input.md> [-o output.pptx]
  python3 md2pptx.py <input.md> --pandoc   # pandoc 経路コマンドを表示

依存:
  pip install python-pptx>=0.6.23 lxml
"""

import argparse
import re
import sys
from pathlib import Path
import lxml.etree as etree

# ============================================================
# カラーパレット（marp-theme.css :root 変数に準拠）
# ============================================================
COLOR_PRIMARY       = "2563eb"   # --color-primary       青
COLOR_PRIMARY_DARK  = "1d4ed8"   # --color-primary-dark  濃青
COLOR_PRIMARY_LIGHT = "3b82f6"   # --color-primary-light 明青
COLOR_ACCENT        = "f59e0b"   # --color-accent        アンバー
COLOR_TEXT          = "1e293b"   # --color-text          濃紺（テキスト）
COLOR_TEXT_LIGHT    = "64748b"   # --color-text-light    グレー
COLOR_BG            = "ffffff"   # --color-bg            白
COLOR_BG_SUB        = "f8fafc"   # --color-bg-sub        薄グレー
COLOR_BORDER        = "e2e8f0"   # --color-border        薄ボーダー
COLOR_WHITE         = "ffffff"   # 白テキスト用
COLOR_PRIMARY_LIGHTER = "eff6ff" # --color-primary-lighter 最薄青
COLOR_ACCENT_LIGHT  = "fffbeb"   # --color-accent-light    薄アンバー
COLOR_CODE_BG       = "1e293b"   # section pre background  暗コード背景
COLOR_CODE_TEXT     = "e2e8f0"   # section pre code color  薄グレーテキスト

# フォント（すべて Meiryo UI で統一）
FONT_NAME = "Meiryo UI"

# フォントサイズ（pt）— Marp CSSの px * 0.9 に調整（視覚的一致）
PT_H1    = 36   # 39px → 36pt（表紙タイトル）
PT_H2    = 26   # 31px → 26pt（スライドタイトル）
PT_H3    = 21   # 25px → 21pt（セクション内見出し）
PT_BODY  = 18   # 20px → 18pt（本文・箇条書き）
PT_SUB   = 15   # 17px → 15pt（図形内・表セル）
PT_NOTE  = 12   # 14px → 12pt（補足）
PT_CAP   = 10   # 12px → 10pt（キャプション）

# スライドサイズ（16:9 ワイドスクリーン）
SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.5

# タイトルバーの高さ（Marp の --title-area-height: 80px 相当）
TITLE_BAR_H_IN = 0.9


# ============================================================
# XML ユーティリティ（CJKフォント・アライメントの確実な設定）
# ============================================================

def _qn(tag: str) -> str:
    """短縮タグ名 → 完全修飾名。例: 'a:latin' → '{http://...}latin'"""
    from pptx.oxml.ns import qn
    return qn(tag)


# OOXML CT_ShapeProperties 子要素の厳密な順序（ECMA-376 準拠）
# PowerPoint は schema 違反の spPr を silently drop するため、
# SubElement（末尾追加）ではなく必ずこの順序で insert すること。
_SPPR_ORDER = [
    "a:xfrm",
    "a:custGeom", "a:prstGeom",
    "a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill",
    "a:ln",
    "a:effectLst", "a:effectDag",
    "a:scene3d",
    "a:sp3d",
    "a:extLst",
]


def _insert_spPr_child(spPr, element):
    """
    spPr の OOXML schema 順序を守って子要素を挿入する。
    同種の既存要素は呼び出し元で事前に削除すること。
    順序が不明な要素は末尾に append する。
    """
    tag_local = etree.QName(element.tag).localname
    prefix_tag = f"a:{tag_local}"
    try:
        my_order = _SPPR_ORDER.index(prefix_tag)
    except ValueError:
        spPr.append(element)
        return
    # 自分より後ろに来るべき最初の兄弟の直前に insert
    for i, child in enumerate(spPr):
        child_local = etree.QName(child.tag).localname
        child_prefix = f"a:{child_local}"
        if child_prefix in _SPPR_ORDER and _SPPR_ORDER.index(child_prefix) > my_order:
            spPr.insert(i, element)
            return
    spPr.append(element)


def _rgb(hex_str: str):
    """'2563eb' → RGBColor(0x25, 0x63, 0xeb)"""
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _clear_shape_effects(shape):
    """
    図形のシャドウ・グロー・反射等のエフェクトをすべて除去する。
    python-pptx の add_shape() はテーマ参照の <p:style> を自動付与し、
    それがシャドウを引き起こす。<p:style> 除去 + <a:effectLst/> 空化で
    純粋なフラット図形にする。
    """
    sp = shape._element

    # テーマ由来のシェイプスタイル参照を削除（シャドウの主要因）
    style_el = sp.find(_qn("p:style"))
    if style_el is not None:
        sp.remove(style_el)

    # spPr 内の effectLst を空にする（schema 順序を守って insert）
    spPr = sp.find(_qn("p:spPr"))
    if spPr is not None:
        for eff in list(spPr.findall(_qn("a:effectLst"))):
            spPr.remove(eff)
        _insert_spPr_child(spPr, etree.Element(_qn("a:effectLst")))  # 空 = エフェクトなし


def _clear_run_effects(run):
    """テキスト run のシャドウ等エフェクトを除去する。"""
    rPr = run._r.get_or_add_rPr()
    for eff in list(rPr.findall(_qn("a:effectLst"))):
        rPr.remove(eff)
    etree.SubElement(rPr, _qn("a:effectLst"))  # 空 = エフェクトなし


def _set_font_cjk(run, name: str, size_pt: float, bold: bool, color_hex: str):
    """
    フォント・サイズ・太字・色を確実に設定する。
    latin / ea（東アジア）/ cs（複合スクリプト）の3要素すべてに
    フォント名を設定することで日本語文字にも確実に適用される。
    （Issue #768 の回避策）
    """
    from pptx.util import Pt

    run.font.size  = Pt(size_pt)
    run.font.bold  = bold
    run.font.color.rgb = _rgb(color_hex)   # 必ず明示（省略で不可視バグ発生）

    # XML 直接操作で latin/ea/cs すべてに設定
    rPr = run._r.get_or_add_rPr()
    for tag in (_qn("a:latin"), _qn("a:ea"), _qn("a:cs")):
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set("typeface", name)

    # テキストシャドウ除去
    _clear_run_effects(run)


def _set_para_align(p, h_align: str = "left"):
    """水平揃えをパラグラフレベルで設定する。"""
    from pptx.enum.text import PP_ALIGN
    mapping = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }
    p.alignment = mapping.get(h_align, PP_ALIGN.LEFT)


def _set_body_anchor(tf, anchor: str = "t"):
    """
    垂直揃えを bodyPr XML 直接操作で設定する。
    anchor: 't'=上揃え, 'ctr'=中央, 'b'=下揃え
    （text_frame.vertical_anchor は信頼性低 — Issue #191 の回避策）
    """
    bodyPr = tf._txBody.find(_qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", anchor)


def _set_rounded_corners(shape, adj: int = 8000):
    """
    図形の角を丸くする（prstGeom を roundRect に変更）。
    adj: 丸みの量（0〜50000、8000≒8px相当）
    CRITICAL: `_insert_spPr_child` で schema 順序を守って挿入する。
    SubElement（末尾追加）は OOXML schema 違反となり PowerPoint が figure を drop する。
    """
    spPr = shape._element.find(_qn("p:spPr"))
    if spPr is None:
        return
    # 既存 prstGeom を削除
    for el in list(spPr.findall(_qn("a:prstGeom"))):
        spPr.remove(el)
    # schema 順序を守って insert（xfrm の直後、Fill の前）
    prstGeom = etree.Element(_qn("a:prstGeom"))
    prstGeom.set("prst", "roundRect")
    avLst = etree.SubElement(prstGeom, _qn("a:avLst"))
    gd = etree.SubElement(avLst, _qn("a:gd"))
    gd.set("name", "adj")
    gd.set("fmla", f"val {adj}")
    _insert_spPr_child(spPr, prstGeom)


def _set_body_inset(tf, left_in=0.1, top_in=0.05, right_in=0.1, bottom_in=0.05):
    """テキストフレームの内側余白を設定する（EMU換算）。"""
    from pptx.util import Inches
    bodyPr = tf._txBody.find(_qn("a:bodyPr"))
    if bodyPr is not None:
        emu = lambda v: str(int(Inches(v)))
        bodyPr.set("lIns", emu(left_in))
        bodyPr.set("tIns", emu(top_in))
        bodyPr.set("rIns", emu(right_in))
        bodyPr.set("bIns", emu(bottom_in))


def _disable_shape_line(shape):
    """図形の枠線を非表示にする。"""
    shape.line.fill.background()


def _add_run(p, text: str, size_pt: float, bold: bool,
             color_hex: str, h_align: str = "left"):
    """段落に run を追加してフォント・色を設定する。"""
    text = text.lstrip("\n")   # 先頭改行バグ（Issue #740）防止
    _set_para_align(p, h_align)
    run = p.add_run()
    run.text = text
    _set_font_cjk(run, FONT_NAME, size_pt, bold, color_hex)
    return run


# ============================================================
# Markdown クリーナ
# ============================================================

def clean_md(text: str) -> str:
    """Markdown 記法をプレーンテキストに変換する。"""
    text = re.sub(r"`([^`]+)`", r"\1", text)          # コードスパン
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)      # 太字 **
    text = re.sub(r"__(.+?)__", r"\1", text)           # 太字 __
    text = re.sub(r"\*(.+?)\*", r"\1", text)           # 斜体 *
    text = re.sub(r"_(.+?)_", r"\1", text)             # 斜体 _
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)  # リンク
    text = re.sub(r"<[^>]+>", "", text)                # HTML タグ
    return text.strip()


# ============================================================
# Marp Markdown パーサ
# ============================================================

def parse_marp_md(md_text: str) -> list[dict]:
    """Marp Markdown をスライドのリストに分解する。"""
    md_text = re.sub(r"^---\n.*?\n---\n", "", md_text, flags=re.DOTALL)
    slide_blocks = re.split(r"\n---\n", md_text)
    slides = []
    for block in slide_blocks:
        block = block.strip()
        if block:
            slides.append(_parse_slide_block(block))
    return slides


def _extract_divs(block: str):
    """
    <div class="NAME">...</div> を探し、クラス名と内側の各 <div> 内容リストを返す。
    貪欲マッチ(.*)で外側 div 全体を正確に捉え、複数の内側 div を全件取得する。
    """
    # 外側 div（class 属性付き）を貪欲マッチで検索
    outer = re.search(
        r'<div\s+class=["\']([^"\']+)["\']>(.*)\n?</div>',
        block, re.DOTALL | re.IGNORECASE
    )
    if not outer:
        return None, []

    layout_class = outer.group(1).split()[0]
    inner = outer.group(2)

    # 内側の <div> ブロックを分割（開始タグで区切る）
    col_parts = re.split(r"<div[^>]*>", inner)
    cols = []
    for part in col_parts:
        # 末尾の </div> を除去してトリム
        part = re.sub(r"\s*</div>\s*$", "", part).strip()
        if part:
            cols.append(part)
    return layout_class, cols


def _extract_code_blocks(text: str) -> list[str]:
    """コードブロック（```...```）の内容を抽出する。"""
    return re.findall(r"```[^\n]*\n(.*?)```", text, re.DOTALL)


def _extract_h3_sections(block: str) -> list[dict]:
    """
    h3 見出しと直後のコンテンツ（箇条書き・コードブロック）を対にして返す。
    戻り値: [{"heading": str, "bullets": [str], "code": str}, ...]
    """
    sections = []
    # h3 の位置でブロックを分割
    parts = re.split(r"^###\s+(.+)$", block, flags=re.MULTILINE)
    # parts[0]: h3 前のテキスト, parts[1]: h3見出し, parts[2]: その後のテキスト, ...
    i = 1
    while i < len(parts) - 1:
        heading = parts[i].strip()
        content = parts[i + 1]

        bullets = re.findall(r"^[-*]\s+(.+)", content, re.MULTILINE)
        code_m = re.search(r"```[^\n]*\n(.*?)```", content, re.DOTALL)
        code = code_m.group(1).strip() if code_m else ""

        sections.append({"heading": heading, "bullets": bullets, "code": code})
        i += 2
    return sections


def _parse_slide_block(block: str) -> dict:
    slide = {
        "class": "",
        "title": "",
        "h1": "",
        "h3_list": [],
        "h3_sections": [],   # [{"heading", "bullets", "code"}]
        "bullets": [],
        "code_blocks": [],   # スライド全体のコードブロック
        "table": [], "table_header": [],
        "layout": "",
        "columns": [],
        "raw_body": "",
        "notes": "",
    }

    # スピーカーノート抽出（_class ディレクティブは除く）
    note_parts = []
    def _extract_note(m):
        c = m.group(1).strip()
        if c.startswith("_"):
            return m.group(0)
        note_parts.append(c)
        return ""
    block = re.sub(r"<!--(.*?)-->", _extract_note, block, flags=re.DOTALL)
    slide["notes"] = "\n".join(note_parts)

    # _class ディレクティブ（コメントが残っている場合）
    m = re.search(r"<!--\s*_class:\s*([^\s>]+)\s*-->", block)
    if m:
        slide["class"] = m.group(1).strip()
        block = block.replace(m.group(0), "").strip()

    # h1（表紙タイトル）
    m = re.match(r"^#\s+(.+)", block, re.MULTILINE)
    if m:
        slide["h1"] = m.group(1).strip()
        block = block.replace(m.group(0), "", 1).strip()

    # h2（スライドタイトル）
    m = re.search(r"^##\s+(.+)", block, re.MULTILINE)
    if m:
        slide["title"] = m.group(1).strip()
        block = block.replace(m.group(0), "", 1).strip()

    # h3 見出し群
    slide["h3_list"] = re.findall(r"^###\s+(.+)", block, re.MULTILINE)

    # h3 + コンテンツのペア（コードブロックや箇条書きを含む）
    slide["h3_sections"] = _extract_h3_sections(block)

    # レイアウト div（貪欲マッチで複数内側 div を全件取得）
    layout_class, cols = _extract_divs(block)
    if layout_class:
        slide["layout"] = layout_class
        slide["columns"] = cols

    # コードブロック
    slide["code_blocks"] = _extract_code_blocks(block)

    # Markdown テーブル
    table_lines = [l.strip() for l in block.splitlines() if re.match(r"^\|", l.strip())]
    if table_lines:
        header, rows = None, []
        for line in table_lines:
            cells = [c.strip() for c in re.split(r"\|", line) if c.strip()]
            if all(re.match(r"^:?-+:?$", c) for c in cells):
                continue
            if header is None:
                header = cells
            else:
                rows.append(cells)
        slide["table_header"] = header or []
        slide["table"] = rows

    # 箇条書き
    slide["bullets"] = re.findall(r"^[-*]\s+(.+)", block, re.MULTILINE)

    slide["raw_body"] = block
    return slide


# ============================================================
# スライド共通ヘルパ
# ============================================================

def _in(v):
    """float インチ → Emu"""
    from pptx.util import Inches
    return Inches(v)


def _pt(v):
    from pptx.util import Pt
    return Pt(v)


def _fill_slide_bg(slide, hex_color: str):
    """スライド全体の背景色を設定する。"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hex_color)


def _add_gradient_rect(slide, x, y, w, h, hex1: str, hex2: str,
                       angle_deg: int = 225):
    """
    グラデーション塗り矩形を追加する。
    angle_deg: CSS の linear-gradient 方向に対応する OOXML 角度（度）
      CSS 135deg（左上→右下）→ OOXML 225 (13500000 * 1/60000)
    """
    shape = slide.shapes.add_shape(1, x, y, w, h)
    _clear_shape_effects(shape)
    _disable_shape_line(shape)

    spPr = shape._element.find(_qn("p:spPr"))
    # 既存のFillを削除してグラデーション追加（schema 順序を守って insert）
    for tag in (_qn("a:solidFill"), _qn("a:gradFill"), _qn("a:noFill")):
        for el in list(spPr.findall(tag)):
            spPr.remove(el)

    gradFill = etree.Element(_qn("a:gradFill"))
    gsLst   = etree.SubElement(gradFill, _qn("a:gsLst"))

    gs1 = etree.SubElement(gsLst, _qn("a:gs"))
    gs1.set("pos", "0")
    c1 = etree.SubElement(gs1, _qn("a:srgbClr"))
    c1.set("val", hex1.lstrip("#"))

    gs2 = etree.SubElement(gsLst, _qn("a:gs"))
    gs2.set("pos", "100000")
    c2 = etree.SubElement(gs2, _qn("a:srgbClr"))
    c2.set("val", hex2.lstrip("#"))

    lin = etree.SubElement(gradFill, _qn("a:lin"))
    lin.set("ang", str(angle_deg * 60000))
    lin.set("scaled", "0")
    # schema 順序を守って spPr に挿入（prstGeom の後、ln の前）
    _insert_spPr_child(spPr, gradFill)
    return shape


def _add_rect(slide, x, y, w, h, fill_hex: str, border_hex: str = None,
              border_pt: float = 0):
    """背景用矩形を追加する（テキストなし）。"""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)
    _clear_shape_effects(shape)          # シャドウ除去
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = _rgb(border_hex)
        shape.line.width = Pt(border_pt)
    else:
        _disable_shape_line(shape)
    return shape


def _add_text_shape(slide, x, y, w, h,
                    text: str, size_pt: float, bold: bool,
                    text_color: str, bg_color: str = None,
                    h_align: str = "left", v_anchor: str = "t",
                    word_wrap: bool = True,
                    border_hex: str = None, border_pt: float = 0):
    """
    テキスト付き図形を追加する。
    フォント・色は必ず明示的に設定する（テーマ依存を排除）。
    """
    from pptx.util import Pt

    shape = slide.shapes.add_shape(1, x, y, w, h)
    _clear_shape_effects(shape)          # シャドウ除去

    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(bg_color)
    else:
        shape.fill.background()

    if border_hex:
        shape.line.color.rgb = _rgb(border_hex)
        shape.line.width = Pt(border_pt)
    else:
        _disable_shape_line(shape)

    tf = shape.text_frame
    tf.word_wrap = word_wrap
    _set_body_anchor(tf, v_anchor)
    _set_body_inset(tf, 0.12, 0.08, 0.12, 0.08)

    p = tf.paragraphs[0]
    text = text.lstrip("\n")
    _add_run(p, text, size_pt, bold, text_color, h_align)
    return shape


def _add_textbox_multi(slide, x, y, w, h,
                       lines: list[tuple],   # [(text, size_pt, bold, color, h_align), ...]
                       v_anchor: str = "t",
                       word_wrap: bool = True):
    """
    複数行テキストボックスを追加する。
    lines の各要素: (text, size_pt, bold, color_hex, h_align)
    """
    shape = slide.shapes.add_shape(1, x, y, w, h)
    _clear_shape_effects(shape)          # シャドウ除去
    shape.fill.background()
    _disable_shape_line(shape)

    tf = shape.text_frame
    tf.word_wrap = word_wrap
    _set_body_anchor(tf, v_anchor)
    _set_body_inset(tf, 0.05, 0.05, 0.05, 0.05)

    first = True
    for (text, size_pt, bold, color_hex, h_align) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        text = text.lstrip("\n")
        _add_run(p, text, size_pt, bold, color_hex, h_align)
    return shape


# ============================================================
# タイトルバー（Marp の fixed h2 帯）
# ============================================================

def _add_title_bar(slide, title: str):
    """
    スライド上部にタイトルバーを追加する。
    CSS: bg-sub背景 + primary-dark テキスト + 下ボーダーのみ（border-bottom: 3px solid primary）。
    """
    W = _in(SLIDE_W_IN)
    H = _in(TITLE_BAR_H_IN)
    BORDER_PT = 2.25  # 3px相当

    # 背景帯（bg-sub 色・ボーダーなし）
    _add_rect(slide, _in(0), _in(0), W, H, fill_hex=COLOR_BG_SUB)

    # 下ボーダーライン（primary 色の細矩形）
    from pptx.util import Pt as _Pt
    border_h = _in(BORDER_PT / 72)  # pt → inch → EMU
    _add_rect(slide, _in(0), H - border_h, W, border_h, fill_hex=COLOR_PRIMARY)

    # タイトルテキスト（primary-dark・太字・左パディング）
    _add_text_shape(
        slide,
        x=_in(0.35), y=_in(0), w=_in(SLIDE_W_IN - 0.5), h=H,
        text=clean_md(title),
        size_pt=PT_H2, bold=True,
        text_color=COLOR_PRIMARY_DARK,
        bg_color=None,
        h_align="left", v_anchor="ctr",
    )


# ============================================================
# スライドレンダラ群
# ============================================================

def _render_lead(slide, data: dict):
    """lead / cta: グラデーション青背景 + 大タイトル中央 + サブタイトル
    CSS: linear-gradient(135deg, #1d4ed8 0%, #2563eb 100%)"""
    _fill_slide_bg(slide, COLOR_PRIMARY_DARK)  # フォールバック（白紙防止）
    # グラデーション全面矩形でオーバーレイ
    _add_gradient_rect(
        slide,
        x=_in(0), y=_in(0), w=_in(SLIDE_W_IN), h=_in(SLIDE_H_IN),
        hex1=COLOR_PRIMARY_DARK, hex2=COLOR_PRIMARY,
        angle_deg=225,  # CSS 135deg = OOXML 225deg
    )

    main_title = clean_md(data["h1"] or data["title"] or "（タイトル）")
    sub_title   = clean_md(data["title"]) if data["h1"] and data["title"] else ""

    # メインタイトル（白・大・中央）
    title_y   = 1.6 if sub_title else 2.3
    title_h   = 2.2 if sub_title else 3.2
    _add_text_shape(
        slide,
        x=_in(0.8), y=_in(title_y),
        w=_in(SLIDE_W_IN - 1.6), h=_in(title_h),
        text=main_title,
        size_pt=PT_H1, bold=True,
        text_color=COLOR_WHITE,
        h_align="center", v_anchor="ctr",
    )

    # サブタイトル（h2 が存在する場合）
    if sub_title:
        _add_text_shape(
            slide,
            x=_in(1.0), y=_in(3.95),
            w=_in(SLIDE_W_IN - 2.0), h=_in(1.0),
            text=sub_title,
            size_pt=PT_H3, bold=False,
            text_color=COLOR_PRIMARY_LIGHT,
            h_align="center", v_anchor="ctr",
        )

    # bullets / plain text（サブタイトルがない場合のフォールバック）
    if not sub_title:
        subs = data["bullets"] or []
        if not subs:
            subs = [l.strip() for l in data["raw_body"].splitlines()
                    if l.strip() and not l.startswith("#")]
        if subs:
            sub_text = "　".join(clean_md(s) for s in subs[:3])
            _add_text_shape(
                slide,
                x=_in(1.2), y=_in(4.8),
                w=_in(SLIDE_W_IN - 2.4), h=_in(1.4),
                text=sub_text,
                size_pt=PT_H3, bold=False,
                text_color=COLOR_BG_SUB,
                h_align="center", v_anchor="ctr",
            )


def _render_section(slide, data: dict):
    """section: bg-sub薄グレー背景 + primaryブルーテキスト中央
    CSS: background-color: var(--color-bg-sub); h2.color: var(--color-primary)"""
    _fill_slide_bg(slide, COLOR_BG_SUB)

    title_text = clean_md(data["title"] or data["h1"] or "（セクション）")

    # タイトル（primary-blue・大・中央）
    _add_text_shape(
        slide,
        x=_in(0.8), y=_in(1.8),
        w=_in(SLIDE_W_IN - 1.6), h=_in(3.0),
        text=title_text,
        size_pt=PT_H1, bold=True,
        text_color=COLOR_PRIMARY,
        h_align="center", v_anchor="ctr",
    )

    # アクセントライン（primary 色）
    _add_rect(slide,
              x=_in(4.5), y=_in(5.0),
              w=_in(4.3), h=_in(0.05),
              fill_hex=COLOR_PRIMARY)


def _render_titled_content(slide, data: dict, accent: bool = False):
    """summary / issue / risk / recommendation: タイトルバー + コンテンツ"""
    title = data["title"]

    # issue/risk は強調バー（primary-dark）
    if accent and title:
        W = _in(SLIDE_W_IN)
        H = _in(TITLE_BAR_H_IN)
        _add_rect(slide, _in(0), _in(0), W, H,
                  fill_hex=COLOR_PRIMARY_DARK)
        _add_text_shape(
            slide,
            x=_in(0.35), y=_in(0), w=_in(SLIDE_W_IN - 0.5), h=H,
            text=clean_md(title),
            size_pt=PT_H2, bold=True,
            text_color=COLOR_WHITE,
            h_align="left", v_anchor="ctr",
        )
    elif title:
        _add_title_bar(slide, title)

    _render_content_area(slide, data, top_in=TITLE_BAR_H_IN + 0.15)


def _render_content_area(slide, data: dict, top_in: float):
    """タイトルバー下のコンテンツを配置する。"""
    avail_h = SLIDE_H_IN - top_in - 0.2

    # テーブルがある場合
    if data["table"] or data["table_header"]:
        _render_table_shapes(slide, data, top_in)
        return

    # h3 + コンテンツのペア（コードブロックや箇条書きが伴う構造化セクション）
    if data["h3_sections"] and any(s["code"] or s["bullets"] for s in data["h3_sections"]):
        _render_h3_sections(slide, data["h3_sections"], top_in, avail_h)
        return

    # 箇条書き
    if data["bullets"]:
        lines = []
        for b in data["bullets"][:8]:
            lines.append(("• " + clean_md(b), PT_BODY, False, COLOR_TEXT, "left"))
        _add_textbox_multi(
            slide,
            x=_in(0.5), y=_in(top_in),
            w=_in(SLIDE_W_IN - 0.8), h=_in(avail_h),
            lines=lines, v_anchor="t",
        )
        return

    # h3 のみ（コンテンツなし）→ 横並びラベル
    if data["h3_list"]:
        _render_h3_row(slide, data["h3_list"], top_in, avail_h)
        return

    # プレーンテキスト
    plain = [l.strip() for l in data["raw_body"].splitlines()
             if l.strip() and not l.startswith("#") and not l.startswith("<")]
    if plain:
        lines = [("• " + clean_md(t), PT_BODY, False, COLOR_TEXT, "left")
                 for t in plain[:8]]
        _add_textbox_multi(
            slide,
            x=_in(0.5), y=_in(top_in),
            w=_in(SLIDE_W_IN - 0.8), h=_in(avail_h),
            lines=lines, v_anchor="t",
        )


def _render_h3_sections(slide, sections: list, top_in: float, avail_h: float):
    """
    h3 見出し + コードブロック/箇条書きを縦積みで配置する。
    例: 使い方（Step 1/2/3）スライド
    """
    n = len(sections)
    if n == 0:
        return

    section_h = avail_h / n
    pad = 0.05

    for i, sec in enumerate(sections):
        y = top_in + i * section_h

        # h3 見出し帯（primary-light 背景）
        label_h = min(0.48, section_h * 0.28)
        _add_text_shape(
            slide,
            x=_in(0.4), y=_in(y + pad),
            w=_in(SLIDE_W_IN - 0.8), h=_in(label_h),
            text=clean_md(sec["heading"]),
            size_pt=PT_H3, bold=True,
            text_color=COLOR_WHITE,
            bg_color=COLOR_PRIMARY,
            h_align="left", v_anchor="ctr",
        )

        content_y = y + pad + label_h + 0.05
        content_h = section_h - label_h - pad * 2 - 0.05

        # コードブロック（CSS: pre { background-color: #1e293b; border-radius: 8px }）
        if sec["code"]:
            code_lines = sec["code"].splitlines()
            display_code = "\n".join(code_lines[:12])
            code_shape = _add_text_shape(
                slide,
                x=_in(0.4), y=_in(content_y),
                w=_in(SLIDE_W_IN - 0.8), h=_in(content_h),
                text=display_code,
                size_pt=PT_NOTE, bold=False,
                text_color=COLOR_CODE_TEXT,   # #e2e8f0 薄グレー
                bg_color=COLOR_CODE_BG,        # #1e293b 暗背景
                border_hex=None, border_pt=0,
                h_align="left", v_anchor="t",
                word_wrap=True,
            )
            _set_rounded_corners(code_shape, adj=6000)  # border-radius: 8px相当
        # 箇条書き
        elif sec["bullets"]:
            lines = [("• " + clean_md(b), PT_SUB, False, COLOR_TEXT, "left")
                     for b in sec["bullets"][:5]]
            _add_textbox_multi(
                slide,
                x=_in(0.5), y=_in(content_y),
                w=_in(SLIDE_W_IN - 0.9), h=_in(content_h),
                lines=lines, v_anchor="t",
            )


def _render_h3_row(slide, h3_list: list, top_in: float, avail_h: float):
    """h3 見出しのみを横並びラベルとして配置する（コンテンツなし）。"""
    n = min(len(h3_list), 3)
    card_w = (SLIDE_W_IN - 0.8) / n
    for i, h3 in enumerate(h3_list[:3]):
        x = 0.4 + i * card_w
        _add_text_shape(
            slide,
            x=_in(x), y=_in(top_in),
            w=_in(card_w - 0.1), h=_in(0.6),
            text=clean_md(h3),
            size_pt=PT_H3, bold=True,
            text_color=COLOR_PRIMARY_DARK,
            h_align="left", v_anchor="ctr",
        )


def _render_kpi(slide, data: dict):
    """kpi: タイトルバー + 4象限 KPI カード"""
    if data["title"]:
        _add_title_bar(slide, data["title"])

    cols = data["columns"] or []
    # class="kpi" の場合、columns が取れていなければ bullets を分割
    if not cols and data["bullets"]:
        cols = data["bullets"]

    grid = [
        (_in(0.35),  _in(TITLE_BAR_H_IN + 0.15)),
        (_in(6.85),  _in(TITLE_BAR_H_IN + 0.15)),
        (_in(0.35),  _in(TITLE_BAR_H_IN + 0.15 + 2.9)),
        (_in(6.85),  _in(TITLE_BAR_H_IN + 0.15 + 2.9)),
    ]
    CARD_W = _in(6.05)
    CARD_H = _in(2.65)

    for idx, (cx, cy) in enumerate(grid):
        # カード背景（bg-sub + border + 角丸）
        card = slide.shapes.add_shape(1, cx, cy, CARD_W, CARD_H)
        _clear_shape_effects(card)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(COLOR_BG_SUB)
        card.line.color.rgb = _rgb(COLOR_BORDER)
        from pptx.util import Pt
        card.line.width = Pt(0.75)
        _set_rounded_corners(card, adj=5000)  # border-radius: 12px相当

        if idx >= len(cols):
            continue

        content = cols[idx]
        text_lines = [clean_md(l.strip("- ").strip())
                      for l in content.splitlines()
                      if l.strip() and not re.match(r"^<", l.strip())]

        tf = card.text_frame
        tf.word_wrap = True
        _set_body_anchor(tf, "ctr")
        _set_body_inset(tf, 0.15, 0.1, 0.15, 0.1)

        first = True
        for j, line in enumerate(text_lines[:4]):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            # 1行目：KPI値として大きく
            if j == 0:
                _add_run(p, line, PT_H2, True, COLOR_PRIMARY, "center")
            else:
                _add_run(p, line, PT_SUB, False, COLOR_TEXT, "center")


def _render_two_column(slide, data: dict, layout_name: str):
    """columns / compare / before-after: 2カラムレイアウト"""
    if data["title"]:
        _add_title_bar(slide, data["title"])

    top = TITLE_BAR_H_IN + 0.15
    avail_h = SLIDE_H_IN - top - 0.2
    is_compare = layout_name in ("compare", "before-after")

    # Before/After ラベル行
    if is_compare:
        label_l = "Before" if layout_name == "before-after" else "現状"
        label_r = "After"  if layout_name == "before-after" else "新方式"
        _add_text_shape(
            slide, x=_in(0.35), y=_in(top),
            w=_in(5.9), h=_in(0.55),
            text=label_l, size_pt=PT_H3, bold=True,
            text_color=COLOR_TEXT_LIGHT,
            h_align="center", v_anchor="ctr",
        )
        _add_text_shape(
            slide, x=_in(7.1), y=_in(top),
            w=_in(5.9), h=_in(0.55),
            text=label_r, size_pt=PT_H3, bold=True,
            text_color=COLOR_PRIMARY,
            h_align="center", v_anchor="ctr",
        )
        top += 0.6
        avail_h -= 0.6

        # 矢印（中央）
        _add_text_shape(
            slide, x=_in(6.25), y=_in(top + avail_h / 2 - 0.3),
            w=_in(0.8), h=_in(0.6),
            text="→", size_pt=PT_H2, bold=True,
            text_color=COLOR_PRIMARY,
            h_align="center", v_anchor="ctr",
        )

    cols = data["columns"]
    left_c  = cols[0] if len(cols) > 0 else ""
    right_c = cols[1] if len(cols) > 1 else ""

    # CSS: .compare > div:first-child { background: #eff6ff } / :last-child { background: #fffbeb }
    if layout_name == "compare":
        left_bg  = COLOR_PRIMARY_LIGHTER   # #eff6ff 薄青
        right_bg = COLOR_ACCENT_LIGHT      # #fffbeb 薄アンバー
    else:
        left_bg  = COLOR_BG_SUB
        right_bg = COLOR_PRIMARY_LIGHTER

    # 左カラム
    _render_column_card(slide, left_c,
                        x_in=0.35, top_in=top,
                        w_in=5.9, h_in=avail_h,
                        bg=left_bg, rounded=True)
    # 右カラム
    _render_column_card(slide, right_c,
                        x_in=7.1, top_in=top,
                        w_in=5.9, h_in=avail_h,
                        bg=right_bg, rounded=True)


def _render_column_card(slide, content: str, x_in, top_in, w_in, h_in,
                        bg: str = COLOR_BG_SUB, rounded: bool = False):
    """カラムカードを描画する。"""
    from pptx.util import Pt

    # カード背景（角丸オプション付き）
    card = slide.shapes.add_shape(
        1, _in(x_in), _in(top_in), _in(w_in), _in(h_in)
    )
    _clear_shape_effects(card)
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(bg)
    card.line.color.rgb = _rgb(COLOR_BORDER)
    card.line.width = Pt(0.75)
    if rounded:
        _set_rounded_corners(card, adj=4000)  # border-radius: 8px相当

    if not content.strip():
        return

    lines_raw = [l.strip() for l in content.splitlines()
                 if l.strip() and not l.startswith("<")]
    if not lines_raw:
        return

    tf = card.text_frame
    tf.word_wrap = True
    _set_body_anchor(tf, "t")
    _set_body_inset(tf, 0.15, 0.12, 0.15, 0.1)

    first = True
    for line in lines_raw[:8]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if re.match(r"^#{1,3}\s", line):
            text = clean_md(re.sub(r"^#{1,3}\s+", "", line))
            _add_run(p, text, PT_H3, True, COLOR_PRIMARY_DARK, "left")
        elif re.match(r"^[-*]\s", line):
            _add_run(p, "• " + clean_md(line[2:]), PT_BODY, False, COLOR_TEXT, "left")
        else:
            _add_run(p, clean_md(line), PT_BODY, False, COLOR_TEXT, "left")


def _render_flow(slide, data: dict):
    """flow: 横並びボックス（primary背景・白テキスト）+ 矢印"""
    if data["title"]:
        _add_title_bar(slide, data["title"])

    items = data["columns"] or []
    if not items:
        items = data["bullets"]

    n = max(len(items), 1)
    margin = 0.35
    total_w = SLIDE_W_IN - margin * 2
    # 矢印幅 0.45" × (n-1) 分を引いてボックス幅を計算
    arrow_w = 0.45
    box_w = (total_w - arrow_w * (n - 1)) / n

    top_in   = TITLE_BAR_H_IN + 0.3
    box_h_in = SLIDE_H_IN - top_in - 0.35

    for i, item in enumerate(items[:5]):
        bx = margin + i * (box_w + arrow_w)

        # CSS: .flow > div { background: bg-sub; border: 1px solid border; border-radius: 8px }
        box = slide.shapes.add_shape(
            1, _in(bx), _in(top_in), _in(box_w), _in(box_h_in)
        )
        _clear_shape_effects(box)
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(COLOR_BG_SUB)
        from pptx.util import Pt as _Pt2
        box.line.color.rgb = _rgb(COLOR_BORDER)
        box.line.width = _Pt2(0.75)
        _set_rounded_corners(box, adj=4000)  # border-radius: 8px

        text_lines = [clean_md(l.strip("- ").strip())
                      for l in item.splitlines()
                      if l.strip() and not re.match(r"^<", l.strip())]

        tf = box.text_frame
        tf.word_wrap = True
        _set_body_anchor(tf, "ctr")
        _set_body_inset(tf, 0.12, 0.1, 0.12, 0.1)

        first = True
        for j, line in enumerate(text_lines[:4]):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            # 1行目はh3サイズ・primary色、以降は本文サイズ・通常色
            if j == 0:
                _add_run(p, line, PT_SUB + 1, True, COLOR_PRIMARY_DARK, "center")
            else:
                _add_run(p, line, PT_SUB, False, COLOR_TEXT, "center")

        # 矢印（最後以外）
        if i < len(items) - 1:
            ax = bx + box_w + 0.03
            _add_text_shape(
                slide,
                x=_in(ax), y=_in(top_in + box_h_in / 2 - 0.28),
                w=_in(arrow_w - 0.06), h=_in(0.56),
                text="→",
                size_pt=PT_H3, bold=True,
                text_color=COLOR_PRIMARY,
                h_align="center", v_anchor="ctr",
            )


def _render_table_shapes(slide, data: dict, top_in: float):
    """テーブルを PPTX table で描画する。"""
    from pptx.util import Pt

    header = data["table_header"]
    rows_data = data["table"]

    n_cols = max(len(header), max((len(r) for r in rows_data), default=1))
    n_header_rows = 1 if header else 0
    n_data_rows = min(len(rows_data), 6)
    n_rows = n_header_rows + n_data_rows

    if n_rows == 0 or n_cols == 0:
        return

    tbl_top = _in(top_in + 0.1)
    tbl_left = _in(0.4)
    tbl_w = _in(SLIDE_W_IN - 0.8)
    avail_h = SLIDE_H_IN - top_in - 0.3
    tbl_h = _in(avail_h)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h)
    tbl = tbl_shape.table

    row_offset = 0
    if header:
        for j, cell_text in enumerate(header[:n_cols]):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(COLOR_PRIMARY)
            p = cell.text_frame.paragraphs[0]
            _set_para_align(p, "center")
            run = p.add_run()
            run.text = clean_md(cell_text)
            _set_font_cjk(run, FONT_NAME, PT_SUB, True, COLOR_WHITE)
        row_offset = 1

    for i, row in enumerate(rows_data[:6]):
        bg = COLOR_BG if i % 2 == 0 else COLOR_BG_SUB
        for j, cell_text in enumerate(row[:n_cols]):
            cell = tbl.cell(i + row_offset, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)
            p = cell.text_frame.paragraphs[0]
            _set_para_align(p, "left")
            run = p.add_run()
            run.text = clean_md(str(cell_text))
            _set_font_cjk(run, FONT_NAME, PT_SUB, False, COLOR_TEXT)


def _render_default(slide, data: dict):
    """デフォルト: タイトルバー + コンテンツ（箇条書き / h3 / テーブル）"""
    if data["title"]:
        _add_title_bar(slide, data["title"])

    top = TITLE_BAR_H_IN + 0.15
    _render_content_area(slide, data, top_in=top)


# ============================================================
# PPTX ビルダ（エントリ）
# ============================================================

def build_pptx(slides: list[dict], output_path: str):
    """スライドリストから PPTX を生成する。"""
    try:
        from pptx import Presentation
    except ImportError:
        print("[ERROR] python-pptx が未インストールです。pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width  = _in(SLIDE_W_IN)
    prs.slide_height = _in(SLIDE_H_IN)

    for slide_data in slides:
        blank_layout = prs.slide_layouts[6]   # Blank
        slide = prs.slides.add_slide(blank_layout)

        # デフォルト背景（白）
        _fill_slide_bg(slide, COLOR_BG)

        cls    = slide_data["class"]
        layout = slide_data["layout"]

        if cls in ("lead", "cta"):
            _render_lead(slide, slide_data)
        elif cls == "section":
            _render_section(slide, slide_data)
        elif cls in ("summary", "issue", "risk", "recommendation"):
            _render_titled_content(slide, slide_data, accent=(cls in ("issue", "risk")))
        elif cls == "kpi" or layout == "kpi":
            _render_kpi(slide, slide_data)
        elif layout in ("columns", "compare", "before-after"):
            _render_two_column(slide, slide_data, layout)
        elif layout == "flow":
            _render_flow(slide, slide_data)
        elif slide_data["table"] or slide_data["table_header"]:
            _render_table_shapes(slide, slide_data, top_in=TITLE_BAR_H_IN + 0.15)
            if slide_data["title"]:
                _add_title_bar(slide, slide_data["title"])
        else:
            _render_default(slide, slide_data)

    prs.save(output_path)
    print(f"[OK] PPTX を生成しました: {output_path}")


# ============================================================
# pandoc フォールバック案内
# ============================================================

def show_pandoc_guide(input_md: str, output_pptx: str):
    print("[INFO] pandoc 経路での PPTX 生成コマンド:")
    print(f'  pandoc "{input_md}" -t pptx -o "{output_pptx}"')
    print()
    print("注意: pandoc 経路はテキスト中心デッキ向けです。")
    print("      columns / kpi / flow 等の高度なレイアウトは python-pptx 経路を推奨します。")


# ============================================================
# エントリポイント
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="Marp Markdown → PPTX 変換スクリプト（aide-pm-slide 用）"
    )
    parser.add_argument("input", help="Marp Markdown ファイルパス（.md）")
    parser.add_argument("-o", "--output",
                        help="出力 PPTX ファイルパス（省略時: 入力と同名の .pptx）")
    parser.add_argument("--pandoc", action="store_true",
                        help="pandoc 経路のコマンドを表示して終了")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"[ERROR] ファイルが見つかりません: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = args.output or str(input_path.with_suffix(".pptx"))

    if args.pandoc:
        show_pandoc_guide(str(input_path), output_path)
        return

    print(f"[INFO] 変換中: {input_path} → {output_path}")
    md_text = input_path.read_text(encoding="utf-8")
    slides  = parse_marp_md(md_text)
    print(f"[INFO] スライド数: {len(slides)}")
    build_pptx(slides, output_path)


if __name__ == "__main__":
    main()
