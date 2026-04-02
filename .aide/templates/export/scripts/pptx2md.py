#!/usr/bin/env python3
"""PowerPoint (.pptx) → Markdown 変換スクリプト

python-pptx でスライドごとにタイトル・本文・ノート・テーブルを抽出し、
Markdown形式で出力する。

Usage:
    python3 pptx2md.py <input.pptx> [--output <output.md>] [--images-dir <dir>]
"""

import argparse
import sys
from pathlib import Path


def extract_table_as_md(table) -> list[str]:
    """pptx テーブルを Markdown テーブルに変換する。"""
    lines: list[str] = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
        lines.append("| " + " | ".join(cells) + " |")
        if i == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return lines


def convert_pptx(input_path: Path, output_path: Path, images_dir: Path) -> dict:
    """pptx を Markdown に変換する。統計情報を返す。"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("エラー: python-pptx が未インストールです。", file=sys.stderr)
        print("  pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(input_path))
    lines: list[str] = []
    stats = {"slides": 0, "images": 0, "tables": 0}
    image_count = 0

    lines.append(f"# {input_path.stem}")
    lines.append("")

    for slide_num, slide in enumerate(prs.slides, 1):
        stats["slides"] += 1

        # タイトル取得
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        lines.append(f"## スライド {slide_num}" + (f": {title}" if title else ""))
        lines.append("")

        # 本文テキストとテーブル・画像を抽出
        for shape in slide.shapes:
            # テーブル
            if shape.has_table:
                stats["tables"] += 1
                lines.extend(extract_table_as_md(shape.table))
                lines.append("")
                continue

            # 画像
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                stats["images"] += 1
                images_dir.mkdir(parents=True, exist_ok=True)
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                img_name = f"slide{slide_num}_image{image_count}.{ext}"
                img_path = images_dir / img_name
                with open(img_path, "wb") as f:
                    f.write(image.blob)
                rel_path = f"{images_dir.name}/{img_name}"
                lines.append(f"![image]({rel_path})")
                lines.append("")
                continue

            # テキストフレーム
            if shape.has_text_frame:
                # タイトル形状はスキップ（既に見出しとして出力済み）
                if shape == slide.shapes.title:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # リスト項目の判定（インデントレベル）
                    level = para.level if para.level else 0
                    if level > 0:
                        indent = "  " * (level - 1)
                        lines.append(f"{indent}- {text}")
                    else:
                        lines.append(text)
                lines.append("")

        # スピーカーノート
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append(f"> **ノート:** {notes_text}")
                lines.append("")

        lines.append("---")
        lines.append("")

    output_path.write_text("\n".join(lines), encoding="utf-8")
    return stats


def main():
    parser = argparse.ArgumentParser(description="PowerPoint (.pptx) → Markdown 変換")
    parser.add_argument("input", help="入力 .pptx ファイルパス")
    parser.add_argument("--output", "-o", help="出力 .md ファイルパス（省略時: 入力と同名.md）")
    parser.add_argument("--images-dir", help="画像抽出先ディレクトリ（省略時: <入力ファイル名>_images）")
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"エラー: ファイルが見つかりません: {input_path}", file=sys.stderr)
        sys.exit(1)
    if input_path.suffix.lower() != ".pptx":
        print(f"エラー: .pptx ファイルを指定してください: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output).resolve() if args.output else input_path.with_suffix(".md")
    images_dir = Path(args.images_dir).resolve() if args.images_dir else input_path.parent / f"{input_path.stem}_images"

    stats = convert_pptx(input_path, output_path, images_dir)

    print(f"変換完了: {input_path.name} → {output_path.name}")
    print(f"  スライド: {stats['slides']}枚  画像: {stats['images']}件  テーブル: {stats['tables']}件")
    if images_dir.exists() and any(images_dir.iterdir()):
        print(f"  画像出力先: {images_dir}")


if __name__ == "__main__":
    main()
